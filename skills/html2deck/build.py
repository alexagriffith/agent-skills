#!/usr/bin/env python3
"""
build.py — assemble sliced slides into a .pptx and/or a .pdf.

Reads slice-report.json + the rendered slide PNGs (produced by slice.py) and
emits deck.pptx and/or deck.pdf, one full-bleed image per 16:9 slide.

Each .pptx slide gets:
  - the unit heading as its TITLE (so an outline/screen-reader sees real text),
  - a full-bleed image of the rendered slide,
  - the unit's prose in the SPEAKER NOTES.

The PDF is Pillow-driven (each PNG becomes one page) — no browser needed, so it
never hangs and matches the pixels a reviewer already approved.

TODO: native editable text boxes/shapes later — image-per-slide for now.

Usage:
    python3 build.py <out_dir> [--pptx] [--pdf] [--name deck]
    (with neither flag, builds both)
"""
import json, argparse
from pathlib import Path

EMU_PER_IN = 914400
DECK_W_IN, DECK_H_IN = 13.333, 7.5      # 16:9 at PowerPoint's widescreen size


def _load(outdir):
    report = json.loads((Path(outdir) / "slice-report.json").read_text())
    slides = report["slides"]
    if not slides:
        raise SystemExit(f"No slides in {outdir}/slice-report.json — run slice.py first.")
    return report, slides


def build_pptx(outdir, slides, path):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(DECK_W_IN)
    prs.slide_height = Inches(DECK_H_IN)
    blank = prs.slide_layouts[6]            # fully blank layout

    for s in slides:
        slide = prs.slides.add_slide(blank)
        img = Path(outdir) / s["png"]
        if not img.exists():
            raise SystemExit(f"Missing slide image: {img}")
        # full-bleed image
        slide.shapes.add_picture(str(img), 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)
        # real title text (off-canvas so it doesn't cover the image but stays
        # in the outline / accessibility tree)
        title = s.get("title") or f"Slide {s['n']}"
        if s.get("part"):
            title = f"{title} ({s['part']})"
        tb = slide.shapes.add_textbox(0, -Inches(1), prs.slide_width, Inches(0.8))
        tb.text_frame.text = title
        # speaker notes: the unit's prose
        notes = s.get("notes_text") or ""
        slide.notes_slide.notes_text_frame.text = notes

    prs.save(str(path))
    return path


def build_pdf(outdir, slides, path):
    from PIL import Image

    pages = []
    for s in slides:
        img = Path(outdir) / s["png"]
        if not img.exists():
            raise SystemExit(f"Missing slide image: {img}")
        im = Image.open(img)
        pages.append(im.convert("RGB") if im.mode != "RGB" else im)
    first, rest = pages[0], pages[1:]
    first.save(str(path), "PDF", save_all=True, append_images=rest, resolution=150.0)
    return path


def main():
    ap = argparse.ArgumentParser(description="Assemble sliced slides into pptx/pdf.")
    ap.add_argument("out_dir", help="the _out dir produced by slice.py")
    ap.add_argument("--pptx", action="store_true")
    ap.add_argument("--pdf", action="store_true")
    ap.add_argument("--name", default="deck", help="output basename (default: deck)")
    args = ap.parse_args()

    do_pptx = args.pptx or not (args.pptx or args.pdf)
    do_pdf = args.pdf or not (args.pptx or args.pdf)

    _report, slides = _load(args.out_dir)
    out = Path(args.out_dir)

    if do_pptx:
        p = build_pptx(args.out_dir, slides, out / f"{args.name}.pptx")
        print(f"wrote {p}  ({len(slides)} slides)")
    if do_pdf:
        p = build_pdf(args.out_dir, slides, out / f"{args.name}.pdf")
        print(f"wrote {p}  ({len(slides)} pages)")


if __name__ == "__main__":
    main()
